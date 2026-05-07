"""W7 후속 — `PptxParser` 단위 테스트 (DE-68 ship).

설계
- python-pptx 의 `Presentation()` 으로 메모리 PPTX 합성 → 파서 직접 호출
- 외부 sample 파일 의존성 0 — fixture 불필요
- 케이스: title placeholder · 텍스트 박스 · 표 · GroupShape 재귀 · 빈 슬라이드

E2 4차 ship — 실 PPTX 자산 fixture 회귀 추가 (`PptxParserRealAssetTest`).
메모리 합성 binary 로는 잡히지 않는 실 파일 회귀를 보호한다. 자산 미존재 시 자동 skip.

자산 디렉토리 우선순위 (5단계, `test_pymupdf_heading.py` 패턴과 정합)
- 1순위: 공개 fixture `<repo>/assets/public/` — 모든 컴퓨터·CI 자동 회귀
- 2순위: `<repo>/assets/` 직속 (사용자 PC raw 자료, `.gitignore` `/assets/*`)
- 3순위: `<repo>/` 루트 직속 (다른 컴퓨터 패턴)
- 4순위: `JETRAG_TEST_PPTX_DIR` ENV 폴백
- 5단계: 자산 부재 시 자동 skip (CI 호환)

stdlib unittest + python-pptx (이미 의존성). 의존성 추가 0.
"""

from __future__ import annotations

import io
import os
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


# repo root 자동 인식: api/tests/test_*.py → parents[2] = repo root
_REPO_ROOT = Path(__file__).resolve().parents[2]
_PUBLIC_PPTX_DIR = _REPO_ROOT / "assets" / "public"

# 공개 fixture (현재 0건 — follow-up 후보, 라이센스 검토 통과 시 추가)
_PUBLIC_PPTX_FILES: list[str] = []

# 비공개 자료 (사용자 PC `assets/` 직속, `.gitignore` 로 다른 컴퓨터엔 부재)
_PRIVATE_PPTX_FILES = ["브랜딩_스튜디오앤드오어.pptx"]

# 회귀 가능 자산 = 공개 + 비공개 (5단계 우선순위로 자동 해석)
_PPTX_FILES = _PUBLIC_PPTX_FILES + _PRIVATE_PPTX_FILES


def _pptx_path(name: str) -> Path:
    """5단계 우선순위로 PPTX fixture 경로 해석. 부재 시 부재 path 반환 (호출부 skipTest).

    1) `<repo>/assets/public/<name>` — 공개 fixture, 모든 컴퓨터·CI 자동
    2) `<repo>/assets/<name>` — 사용자 PC raw 자료 (`.gitignore` `/assets/*`)
    3) `<repo>/<name>` — 다른 컴퓨터에서 자료가 repo 루트 직속에 있을 때
    4) `$JETRAG_TEST_PPTX_DIR/<name>` — 외장 디스크·별 위치 보강용 ENV 폴백
    5) 부재 → public path 반환 (exists() False, 호출부 skipTest)
    """
    public = _PUBLIC_PPTX_DIR / name
    if public.exists():
        return public

    assets_direct = _REPO_ROOT / "assets" / name
    if assets_direct.exists():
        return assets_direct

    repo_root_direct = _REPO_ROOT / name
    if repo_root_direct.exists():
        return repo_root_direct

    env_base = os.environ.get("JETRAG_TEST_PPTX_DIR")
    if env_base:
        env_path = Path(env_base) / name
        if env_path.exists():
            return env_path

    return public  # exists() False — 호출부에서 skipTest


def _make_pptx_bytes(build_fn) -> bytes:
    """build_fn(prs) 콜백으로 PPTX 합성 → bytes 반환."""
    prs = Presentation()
    build_fn(prs)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class PptxBasicTextTest(unittest.TestCase):
    """텍스트 박스만 있는 단순 슬라이드 — 1 section / page=1."""

    def test_single_slide_with_text(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        def build(prs):
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)
            tx = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(6), Inches(2)
            )
            tf = tx.text_frame
            tf.text = "프레젠테이션 제목"
            tf.add_paragraph().text = "본문 텍스트입니다."

        data = _make_pptx_bytes(build)
        parser = PptxParser()
        result = parser.parse(data, file_name="test.pptx")

        self.assertEqual(len(result.sections), 1)
        self.assertEqual(result.sections[0].page, 1)
        self.assertIn("프레젠테이션 제목", result.sections[0].text)
        self.assertIn("본문 텍스트", result.sections[0].text)
        # title 우선 — 첫 텍스트 박스 의 첫 줄
        self.assertEqual(
            result.sections[0].section_title, "프레젠테이션 제목"
        )


class PptxTitlePlaceholderTest(unittest.TestCase):
    """title placeholder 가 있는 레이아웃 — title 우선 사용."""

    def test_title_placeholder_priority(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        def build(prs):
            title_layout = prs.slide_layouts[0]  # title slide
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = "공식 제목"
            # subtitle/body placeholder 채움
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    ph.text = "부제목 텍스트"
                    break

        data = _make_pptx_bytes(build)
        parser = PptxParser()
        result = parser.parse(data, file_name="title.pptx")

        self.assertEqual(len(result.sections), 1)
        self.assertEqual(result.sections[0].section_title, "공식 제목")


class PptxTableTest(unittest.TestCase):
    """표가 있는 슬라이드 — DocxParser 패턴 (` | ` separator)."""

    def test_table_renders_with_pipe_separator(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        def build(prs):
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            shape = slide.shapes.add_table(
                rows=2, cols=2, left=Inches(1), top=Inches(1),
                width=Inches(4), height=Inches(2),
            )
            tbl = shape.table
            tbl.cell(0, 0).text = "헤더1"
            tbl.cell(0, 1).text = "헤더2"
            tbl.cell(1, 0).text = "값1"
            tbl.cell(1, 1).text = "값2"

        data = _make_pptx_bytes(build)
        parser = PptxParser()
        result = parser.parse(data, file_name="table.pptx")

        self.assertEqual(len(result.sections), 1)
        text = result.sections[0].text
        self.assertIn("헤더1 | 헤더2", text)
        self.assertIn("값1 | 값2", text)


class PptxMultiSlideTest(unittest.TestCase):
    """여러 슬라이드 — page 1·2·3 자동 부여."""

    def test_multi_slide_pages_increment(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        def build(prs):
            blank = prs.slide_layouts[6]
            for i in range(3):
                slide = prs.slides.add_slide(blank)
                tx = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(6), Inches(1)
                )
                tx.text_frame.text = f"슬라이드 {i + 1} 본문"

        data = _make_pptx_bytes(build)
        parser = PptxParser()
        result = parser.parse(data, file_name="multi.pptx")

        self.assertEqual(len(result.sections), 3)
        for i, sec in enumerate(result.sections):
            self.assertEqual(sec.page, i + 1)
            self.assertIn(f"슬라이드 {i + 1}", sec.text)


class PptxEmptySlideTest(unittest.TestCase):
    """텍스트 0인 슬라이드 — section 미생성, raise 0 (graceful)."""

    def test_empty_slide_skipped(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        def build(prs):
            blank = prs.slide_layouts[6]
            prs.slides.add_slide(blank)  # 빈 슬라이드 1개
            slide2 = prs.slides.add_slide(blank)
            slide2.shapes.add_textbox(
                Inches(1), Inches(1), Inches(6), Inches(1)
            ).text_frame.text = "두 번째 슬라이드"

        data = _make_pptx_bytes(build)
        parser = PptxParser()
        result = parser.parse(data, file_name="empty.pptx")

        # 빈 슬라이드는 section 미생성, 두 번째 슬라이드만 1 section
        self.assertEqual(len(result.sections), 1)
        self.assertIn("두 번째 슬라이드", result.sections[0].text)
        # page 는 슬라이드 인덱스 기준이라 두 번째 슬라이드는 page=2
        self.assertEqual(result.sections[0].page, 2)


class PptxCorruptedTest(unittest.TestCase):
    """깨진 PPTX → RuntimeError wrap (DocxParser 패턴 일치)."""

    def test_corrupted_raises_runtime(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        with self.assertRaises(RuntimeError) as ctx:
            PptxParser().parse(b"not a pptx file", file_name="bad.pptx")
        self.assertIn("PPTX 파서 초기화 실패", str(ctx.exception))


class PptxCanParseTest(unittest.TestCase):
    """`can_parse` 가 .pptx 만 True."""

    def test_can_parse_extension_only(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser

        p = PptxParser()
        self.assertTrue(p.can_parse("a.pptx", None))
        self.assertTrue(p.can_parse("/path/to/file.PPTX", None))
        self.assertFalse(p.can_parse("a.docx", None))
        self.assertFalse(p.can_parse("a.pdf", None))


class PptxVisionReroutingTest(unittest.TestCase):
    """W8 Day 2 — 슬라이드 텍스트 0 + Picture 있음 → image_parser.parse() 위임.

    mock ImageParser 로 Vision API 호출 차단 — 외부 의존성 0.
    """

    def _make_pptx_with_picture(self, picture_count: int = 1) -> bytes:
        """슬라이드 1장에 Picture N개 (텍스트 박스 0). PIL 합성 PNG 사용."""
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        # 단색 PNG 합성 (테스트용 — 실 OCR 결과는 mock 으로 대체)
        png_buf = BytesIO()
        Image.new("RGB", (200, 100), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        for i in range(picture_count):
            slide.shapes.add_picture(
                BytesIO(png_bytes),
                left=Inches(1 + i),
                top=Inches(1),
                width=Inches(2),
                height=Inches(1.5),
            )
        out = BytesIO()
        prs.save(out)
        return out.getvalue()

    def test_picture_only_slide_invokes_image_parser(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser
        from app.adapters.parser import ExtractedSection, ExtractionResult

        data = self._make_pptx_with_picture(picture_count=1)

        # mock ImageParser — parse() 호출 횟수 + blob 인자 검증
        parse_calls: list[bytes] = []

        class _FakeImageParser:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None) -> ExtractionResult:
                parse_calls.append(blob)
                return ExtractionResult(
                    source_type="image",
                    sections=[
                        ExtractedSection(
                            text="[표지] 모의 OCR 결과",
                            page=None,
                            section_title="이미지 분류: 표지",
                        ),
                        ExtractedSection(
                            text="모의 OCR 결과", page=None, section_title="OCR 텍스트"
                        ),
                    ],
                    raw_text="[표지] 모의 OCR 결과\n\n모의 OCR 결과",
                    warnings=[],
                )

        parser = PptxParser(image_parser=_FakeImageParser())
        result = parser.parse(data, file_name="picture.pptx")

        self.assertEqual(len(parse_calls), 1, "Vision OCR 1회 호출")
        self.assertGreater(len(parse_calls[0]), 100, "image blob 전달")
        self.assertEqual(len(result.sections), 1)
        self.assertIn("모의 OCR 결과", result.sections[0].text)
        self.assertEqual(result.sections[0].page, 1)
        self.assertEqual(result.sections[0].section_title, "p.1 (Vision OCR)")

    def test_text_slide_skips_image_parser(self) -> None:
        from app.adapters.impl.pptx_parser import PptxParser
        from pptx import Presentation
        from pptx.util import Inches

        # 텍스트 박스만 있는 슬라이드 — Vision 호출 0회 기대
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(6), Inches(2)
        )
        tx.text_frame.text = "텍스트가 있는 슬라이드"
        buf = io.BytesIO()
        prs.save(buf)

        parse_calls: list[bytes] = []

        class _SpyImageParser:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None):
                parse_calls.append(blob)
                raise AssertionError("텍스트 슬라이드는 Vision 호출 X")

        parser = PptxParser(image_parser=_SpyImageParser())
        result = parser.parse(buf.getvalue(), file_name="text.pptx")

        self.assertEqual(parse_calls, [], "텍스트 슬라이드는 Vision 호출 0회")
        self.assertEqual(len(result.sections), 1)
        self.assertIn("텍스트가 있는 슬라이드", result.sections[0].text)

    def test_max_5_slides_cap(self) -> None:
        """6 픽처 슬라이드 → 첫 5개만 Vision rerouting (cap)."""
        from app.adapters.impl.pptx_parser import PptxParser
        from app.adapters.parser import ExtractedSection, ExtractionResult
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        png_buf = BytesIO()
        Image.new("RGB", (100, 50), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        prs = Presentation()
        for _ in range(6):  # 6 슬라이드 (cap=5 검증)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                BytesIO(png_bytes),
                Inches(1), Inches(1), Inches(2), Inches(1),
            )
        buf = BytesIO()
        prs.save(buf)

        parse_calls: list[int] = []

        class _Counting:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None) -> ExtractionResult:
                parse_calls.append(1)
                return ExtractionResult(
                    source_type="image",
                    sections=[
                        ExtractedSection(text="[표지] OCR", page=None, section_title="t")
                    ],
                    raw_text="[표지] OCR",
                    warnings=[],
                )

        parser = PptxParser(image_parser=_Counting())
        result = parser.parse(buf.getvalue(), file_name="cap.pptx")

        self.assertEqual(len(parse_calls), 5, "max 5 cap")
        self.assertEqual(len(result.sections), 5, "6번째는 텍스트도 OCR도 없어 skip")

    def test_vision_failure_graceful(self) -> None:
        """Vision API 실패 시 RuntimeError 흡수, 슬라이드는 skip + warnings."""
        from app.adapters.impl.pptx_parser import PptxParser

        data = self._make_pptx_with_picture(picture_count=1)

        class _BrokenImageParser:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None):
                raise RuntimeError("Gemini API down")

        parser = PptxParser(image_parser=_BrokenImageParser())
        result = parser.parse(data, file_name="broken.pptx")

        self.assertEqual(result.sections, [], "OCR 실패 슬라이드는 section 미생성")
        self.assertTrue(
            any("Vision OCR 실패" in w for w in result.warnings),
            f"warnings 에 graceful 메시지 — got {result.warnings}",
        )

    def test_failure_respects_cap_quota_protection(self) -> None:
        """W9 Day 3 한계 #47 회수 — Vision 실패 시에도 cap (시도 기준) 적용.

        이전 버그: cap 이 *성공* 만 카운트 → 11 slides 모두 호출 → quota 초과 누적.
        fix 후: cap 이 *시도* 기준 → quota 보호 의도 보존.
        """
        from app.adapters.impl.pptx_parser import PptxParser
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        png_buf = BytesIO()
        Image.new("RGB", (100, 50), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        # 11 슬라이드 모두 picture-only — 모든 Vision 호출이 실패하더라도 cap 5 적용
        prs = Presentation()
        for _ in range(11):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                BytesIO(png_bytes),
                Inches(1), Inches(1), Inches(2), Inches(1),
            )
        buf = BytesIO()
        prs.save(buf)

        parse_calls: list[int] = []

        # 메시지에 quota 키워드 미포함 — fast-fail (한계 #49) 분기 회피하고 cap 만 검증
        class _AlwaysFail:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None):
                parse_calls.append(1)
                raise RuntimeError("Service temporarily unavailable")

        parser = PptxParser(image_parser=_AlwaysFail())
        result = parser.parse(buf.getvalue(), file_name="all_fail.pptx")

        self.assertEqual(
            len(parse_calls), 5,
            f"실패 시에도 cap 5 적용 (시도 기준) — got {len(parse_calls)} (이전 버그: 11)",
        )
        self.assertEqual(result.sections, [], "모두 실패라 section 0")

    def test_quota_exhausted_fast_fail(self) -> None:
        """W9 Day 4 한계 #49 — RESOURCE_EXHAUSTED 메시지 시 첫 호출 후 즉시 skip.

        Day 3 cap (5회) 위에 추가 보호 레이어. 11 슬라이드 모두 picture-only +
        ImageParser 가 'RESOURCE_EXHAUSTED' 메시지로 raise → 1회 호출 후 stop.
        """
        from app.adapters.impl.pptx_parser import PptxParser
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        png_buf = BytesIO()
        Image.new("RGB", (100, 50), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        prs = Presentation()
        for _ in range(11):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                BytesIO(png_bytes),
                Inches(1), Inches(1), Inches(2), Inches(1),
            )
        buf = BytesIO()
        prs.save(buf)

        parse_calls: list[int] = []

        class _QuotaExhausted:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None):
                parse_calls.append(1)
                # Gemini SDK 의 google.api_core.exceptions.ResourceExhausted str 형식 모방
                raise RuntimeError(
                    "429 RESOURCE_EXHAUSTED. {'error': {'code': 429, "
                    "'message': 'You exceeded your current quota'}}"
                )

        parser = PptxParser(image_parser=_QuotaExhausted())
        result = parser.parse(buf.getvalue(), file_name="quota.pptx")

        self.assertEqual(
            len(parse_calls), 1,
            f"quota 감지 즉시 fast-fail — got {len(parse_calls)} (cap 5 미적용까지 stop)",
        )
        self.assertEqual(result.sections, [])
        # warnings 에 quota 감지 + 이후 skip 메시지
        self.assertTrue(
            any("quota 감지" in w for w in result.warnings),
            f"warnings 에 quota 메시지 — got {result.warnings}",
        )

    def test_429_in_message_triggers_fast_fail(self) -> None:
        """RESOURCE_EXHAUSTED 가 아닌 다른 형식의 quota 메시지도 감지 — '429' 키워드만으로."""
        from app.adapters.impl.pptx_parser import PptxParser

        data = self._make_pptx_with_picture(picture_count=1)

        parse_calls: list[int] = []

        class _Code429:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None):
                parse_calls.append(1)
                raise RuntimeError("HTTP 429 Too Many Requests from upstream")

        parser = PptxParser(image_parser=_Code429())
        result = parser.parse(data, file_name="429.pptx")

        self.assertEqual(len(parse_calls), 1)
        self.assertEqual(result.sections, [])
        self.assertTrue(any("quota 감지" in w for w in result.warnings))

    def test_no_image_parser_disables_vision(self) -> None:
        """image_parser=None (기본) → Vision 비활성, picture-only 슬라이드는 skip."""
        from app.adapters.impl.pptx_parser import PptxParser

        data = self._make_pptx_with_picture(picture_count=1)
        parser = PptxParser(image_parser=None)
        result = parser.parse(data, file_name="no_vision.pptx")
        self.assertEqual(result.sections, [])


class PptxVisionAugmentTest(unittest.TestCase):
    """W9 Day 1 — 한계 #28 회수: 짧은 텍스트 + Picture 슬라이드 OCR 결합 (augment mode)."""

    def _make_pptx(self, text: str) -> bytes:
        """텍스트 박스 1개 + Picture 1개 슬라이드 합성."""
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        png_buf = BytesIO()
        Image.new("RGB", (200, 100), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 텍스트 박스
        tx = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(6), Inches(1)
        )
        tx.text_frame.text = text
        # Picture
        slide.shapes.add_picture(
            BytesIO(png_bytes),
            Inches(1), Inches(2), Inches(4), Inches(3),
        )
        out = BytesIO()
        prs.save(out)
        return out.getvalue()

    def _make_ocr_image_parser(self, ocr_text: str = "OCR로 회수한 본문 텍스트입니다."):
        """OCR 결과를 mock — ExtractionResult 1 section."""
        from app.adapters.parser import ExtractedSection, ExtractionResult

        parse_calls: list[bytes] = []

        class _Mock:
            def parse(self, blob: bytes, *, file_name: str, source_type: str | None = None) -> ExtractionResult:
                parse_calls.append(blob)
                return ExtractionResult(
                    source_type="image",
                    sections=[
                        ExtractedSection(
                            text=ocr_text, page=None, section_title="OCR"
                        )
                    ],
                    raw_text=ocr_text,
                    warnings=[],
                )

        return _Mock(), parse_calls

    def test_short_text_triggers_augment(self) -> None:
        """텍스트 < 50자 + Picture 있음 → OCR 결합 (텍스트 + OCR 모두 보존)."""
        from app.adapters.impl.pptx_parser import PptxParser

        # 18자 — threshold 50 미만
        data = self._make_pptx("디자인 컨셉 표지 슬라이드")
        mock_ip, parse_calls = self._make_ocr_image_parser(
            ocr_text="OCR 회수: 본문 추가 정보"
        )

        parser = PptxParser(image_parser=mock_ip)
        result = parser.parse(data, file_name="short.pptx")

        self.assertEqual(len(parse_calls), 1, "augment 모드에서 OCR 1회 호출")
        self.assertEqual(len(result.sections), 1)
        text = result.sections[0].text
        # 기존 텍스트 + OCR 모두 포함
        self.assertIn("디자인 컨셉 표지 슬라이드", text)
        self.assertIn("OCR 회수: 본문 추가 정보", text)

    def test_long_text_skips_augment(self) -> None:
        """텍스트 ≥ 50자 → OCR skip (RPD 절약)."""
        from app.adapters.impl.pptx_parser import PptxParser

        long_text = (
            "이 슬라이드는 충분히 긴 본문 텍스트를 가지고 있어서 image OCR 추가 결합이 "
            "필요하지 않은 케이스입니다 — Vision RPD 절약 정책."
        )
        self.assertGreaterEqual(len(long_text), 50)
        data = self._make_pptx(long_text)
        mock_ip, parse_calls = self._make_ocr_image_parser()

        parser = PptxParser(image_parser=mock_ip)
        result = parser.parse(data, file_name="long.pptx")

        self.assertEqual(parse_calls, [], "텍스트 풍부 슬라이드는 OCR 호출 0회")
        self.assertEqual(len(result.sections), 1)
        self.assertIn("충분히 긴 본문", result.sections[0].text)

    def test_augment_respects_max_cap(self) -> None:
        """6 슬라이드 모두 짧은 텍스트 + Picture → 첫 5개만 augment (cap)."""
        from app.adapters.impl.pptx_parser import PptxParser
        from io import BytesIO
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        png_buf = BytesIO()
        Image.new("RGB", (100, 50), color="white").save(png_buf, format="PNG")
        png_bytes = png_buf.getvalue()

        prs = Presentation()
        for i in range(6):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(6), Inches(1)
            )
            tx.text_frame.text = f"짧은 제목 {i + 1}"  # 항상 < 50
            slide.shapes.add_picture(
                BytesIO(png_bytes),
                Inches(1), Inches(2), Inches(2), Inches(1),
            )
        buf = BytesIO()
        prs.save(buf)

        mock_ip, parse_calls = self._make_ocr_image_parser()
        parser = PptxParser(image_parser=mock_ip)
        result = parser.parse(buf.getvalue(), file_name="cap.pptx")

        self.assertEqual(len(parse_calls), 5, "augment 도 max 5 cap 적용")
        self.assertEqual(len(result.sections), 6, "6 슬라이드 모두 section 생성 (텍스트만)")


class PptxParserRealAssetTest(unittest.TestCase):
    """E2 4차 ship — 실 PPTX 자산에 대한 회귀 보호.

    메모리 합성 (`_make_pptx_bytes`) 으로는 잡히지 않는 케이스:
    - 실제 디자인 도구로 생성된 슬라이드 마스터·레이아웃 변형
    - 한국어 글꼴·복합 단락·SmartArt 등 외부 도구 산출물
    - 실 사용자 자산의 인코딩·메타데이터 정합

    자산 부재 시 자동 skip (CI 호환). 사용자 PC `assets/` 직속에 자료 있으면 자동 진입.
    """

    def setUp(self) -> None:
        self._target_name: str | None = None
        self._target_path: Path | None = None
        for name in _PPTX_FILES:
            path = _pptx_path(name)
            if path.exists():
                self._target_name = name
                self._target_path = path
                return
        self.skipTest("실 PPTX 자산 부재 — public/private 모두 미존재")

    def test_can_parse_real_pptx(self) -> None:
        """실 PPTX 파싱이 raise 없이 정상 종료 + ExtractionResult 정합 + 텍스트/메타데이터 검증.

        텍스트 슬라이드가 있는 자산: sections >= 1 + text non-empty + page 정수 검증.
        picture-only 자산 (디자인 PPTX 등): image_parser=None 기본 동작으로 sections=0 정상.
        둘 다 ExtractionResult schema 정합·raise 0 회귀 보호.
        """
        from app.adapters.impl.pptx_parser import PptxParser

        assert self._target_path is not None and self._target_name is not None
        data = self._target_path.read_bytes()
        result = PptxParser().parse(data, file_name=self._target_name)

        # 정합 — 모든 케이스 공통 (raise 0 + ExtractionResult schema)
        self.assertEqual(result.source_type, "pptx")
        self.assertIsInstance(result.sections, list)
        self.assertIsInstance(result.raw_text, str)
        self.assertIsInstance(result.warnings, list)

        # 텍스트 슬라이드 자산은 추가 검증 (picture-only 자산이면 아래는 skip)
        if result.sections:
            self.assertTrue(
                result.sections[0].text.strip(),
                "section[0].text 가 비어 있음 — 파서 텍스트 회수 실패",
            )
            for sec in result.sections:
                self.assertIsInstance(
                    sec.page, int,
                    f"page 가 정수가 아님: {sec.page!r}",
                )
                self.assertGreaterEqual(
                    sec.page, 1,
                    f"page 가 1 미만: {sec.page}",
                )


if __name__ == "__main__":
    unittest.main()
