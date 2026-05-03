"""W7 후속 — `_input_gate.validate_magic` 단위 테스트.

배경
- W7 Day 6 마감 후 사용자 docx/pptx 업로드 시 filetype 1.2.0 가 deep MIME
  (`application/vnd.openxmlformats-officedocument...`) 을 반환하는 케이스 발견.
- 기존 `_EXT_TO_MIMES` 는 `application/zip` 만 허용 → 거부.
- fix 후 (deep MIME 도 허용) 회귀 보호.

stdlib unittest + python-docx/python-pptx (이미 의존성). 의존성 추가 0.
"""

from __future__ import annotations

import io
import unittest

from fastapi import HTTPException


def _make_docx_head_bytes() -> bytes:
    """python-docx 로 빈 DOCX 합성 후 앞 4096 byte 추출."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("test")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()[:4096]


def _make_pptx_head_bytes() -> bytes:
    """python-pptx 로 빈 PPTX 합성 후 앞 4096 byte 추출."""
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()[:4096]


class DocxValidateTest(unittest.TestCase):
    """W7 후속 fix — DOCX deep MIME 도 통과."""

    def test_docx_validates(self) -> None:
        from app.routers._input_gate import validate_magic

        head = _make_docx_head_bytes()
        # raise 없이 통과
        validate_magic(ext=".docx", raw_head=head)


class PptxValidateTest(unittest.TestCase):
    """W7 후속 fix — PPTX deep MIME 도 통과."""

    def test_pptx_validates(self) -> None:
        from app.routers._input_gate import validate_magic

        head = _make_pptx_head_bytes()
        validate_magic(ext=".pptx", raw_head=head)


class ExeMasqueradeTest(unittest.TestCase):
    """exe 가 .docx 위장 — 차단."""

    def test_exe_masquerade_rejected(self) -> None:
        from app.routers._input_gate import validate_magic

        # PE/COFF 시그니처 (Windows exe)
        exe_head = b"MZ\x90\x00" + b"\x00" * 4092
        with self.assertRaises(HTTPException) as ctx:
            validate_magic(ext=".docx", raw_head=exe_head)
        self.assertEqual(ctx.exception.status_code, 400)


class UnsupportedExtTest(unittest.TestCase):
    """매핑 없는 확장자 — 화이트리스트 외 케이스."""

    def test_unknown_ext_rejected(self) -> None:
        from app.routers._input_gate import validate_magic

        with self.assertRaises(HTTPException) as ctx:
            validate_magic(ext=".exe", raw_head=b"\x00" * 4096)
        self.assertEqual(ctx.exception.status_code, 400)


if __name__ == "__main__":
    unittest.main()
