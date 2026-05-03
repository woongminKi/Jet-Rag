"""python-pptx 기반 PPTX 프레젠테이션 파서 (DE-68).

배경
- W4-Q-9 sniff 결과 python-pptx 1.x 의 `Presentation.slides` + `shape.has_text_frame` 으로 충분 판정.
- 페르소나 A 의 회의 발표·교육 자료 빈도 ↑ → 사용자 자료 업로드 시점에 ship (W7 후속, W8 Day 1).

설계
- 슬라이드 1개 = `ExtractedSection` 1개 (raw_text 는 슬라이드 내부 모든 텍스트 join)
- `page` = slide_index + 1 (1-based, search UI 의 "p.N" 표기와 호환)
- `section_title` = title placeholder 우선 → 첫 번째 텍스트 shape fallback
- 표 (Shape with `has_table`): DocxParser `_table_to_text` 와 동일 ` | ` separator
- GroupShape: `.shapes` 재귀 — 디자인 PPT 흔한 그룹 구조

W8 Day 2 — Vision OCR rerouting (한계 #23 회수)
- 슬라이드 텍스트 0 이고 Picture 가 있으면 가장 큰 Picture 의 image_bytes 를 ImageParser 에 위임
- max 5 슬라이드 cap (PyMuPDF 스캔 PDF rerouting 의 `_MAX_SCAN_PAGES` 패턴 일치, Vision RPD 20 제약)
- ImageParser composition — 직접 Gemini 호출 X, 모든 경로 통일
- 디자인 카탈로그 (Picture 100%) 같은 PPTX 자료의 텍스트 회수

graceful degrade
- python-pptx 가 corrupted PPTX 에서 raise → RuntimeError wrap (DocxParser 패턴)
- 슬라이드/shape/Vision 단위 부분 실패는 warnings 누적, 다음 슬라이드 계속
"""

from __future__ import annotations

import io
import logging
from pathlib import PurePosixPath

from pptx import Presentation

from app.adapters.parser import ExtractedSection, ExtractionResult

logger = logging.getLogger(__name__)

# Vision OCR rerouting cap — Gemini Flash RPD 20 + 디자인 PPT 첫 5장이 보통 표지·요약
_MAX_VISION_SLIDES = 5

# 텍스트 풍부도 임계값 (W9 Day 1 — 한계 #28).
# 슬라이드 텍스트가 본 임계값 미만이면 image OCR 추가 결합 — 짧은 제목 + 큰 이미지
# 디자인 PPT 케이스 회수. 일반 텍스트 풍부 슬라이드는 OCR skip 으로 RPD 절약.
_VISION_AUGMENT_TEXT_THRESHOLD = 50


class PptxParser:
    source_type = "pptx"

    def __init__(self, image_parser=None) -> None:
        """`image_parser` 는 텍스트 0 슬라이드의 가장 큰 Picture rerouting 용 (W8 Day 2).

        None 이면 Vision OCR 기능 비활성 — 단위 테스트 / 외부 의존성 회피 케이스.
        extract.py 가 production 에서 `ImageParser()` 인스턴스 주입.
        """
        self._image_parser = image_parser

    def can_parse(self, file_name: str, mime_type: str | None) -> bool:
        ext = PurePosixPath(file_name).suffix.lower()
        return ext == ".pptx"

    def parse(self, data: bytes, *, file_name: str) -> ExtractionResult:
        sections: list[ExtractedSection] = []
        warnings: list[str] = []
        raw_parts: list[str] = []

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(
                f"PPTX 파서 초기화 실패: {file_name}: {exc}"
            ) from exc

        vision_slides_used = 0

        for slide_idx, slide in enumerate(prs.slides):
            try:
                slide_title = _extract_slide_title(slide)
                slide_text_parts = _extract_slide_text(slide, warnings=warnings)

                # W8 Day 2 — 텍스트 부재 시 Vision rerouting (rerouting mode)
                # W9 Day 1 — 텍스트 짧으면 Vision augment (한계 #28, augment mode)
                #   · rerouting: 텍스트 0 → OCR 만 사용
                #   · augment: 텍스트 < 50자 → 기존 텍스트 + OCR 결합
                current_text_len = sum(len(p) for p in slide_text_parts)
                needs_ocr = (
                    self._image_parser is not None
                    and vision_slides_used < _MAX_VISION_SLIDES
                    and current_text_len < _VISION_AUGMENT_TEXT_THRESHOLD
                )
                if needs_ocr:
                    ocr_text = _vision_ocr_largest_picture(
                        slide,
                        slide_idx=slide_idx,
                        file_name=file_name,
                        image_parser=self._image_parser,
                        warnings=warnings,
                    )
                    if ocr_text:
                        vision_slides_used += 1
                        if not slide_text_parts:
                            # rerouting mode — 텍스트 0
                            if not slide_title:
                                slide_title = f"p.{slide_idx + 1} (Vision OCR)"
                            slide_text_parts = [ocr_text]
                        else:
                            # augment mode — 기존 텍스트 + OCR 결합
                            slide_text_parts = [*slide_text_parts, ocr_text]

                if not slide_text_parts:
                    continue

                slide_text = "\n".join(slide_text_parts)
                sections.append(
                    ExtractedSection(
                        text=slide_text,
                        page=slide_idx + 1,
                        section_title=slide_title,
                        bbox=None,
                    )
                )
                raw_parts.append(slide_text)
            except Exception as exc:  # noqa: BLE001 — 슬라이드 단위 부분 실패 허용
                msg = f"PPTX slide {slide_idx + 1} 추출 실패: {exc}"
                warnings.append(msg)
                logger.warning("%s (file=%s)", msg, file_name)
                continue

        if vision_slides_used > 0:
            logger.info(
                "PPTX Vision OCR rerouting: %d slides 처리 (file=%s, cap=%d)",
                vision_slides_used,
                file_name,
                _MAX_VISION_SLIDES,
            )

        return ExtractionResult(
            source_type=self.source_type,
            sections=sections,
            raw_text="\n\n".join(raw_parts),
            warnings=warnings,
        )


def _extract_slide_title(slide) -> str | None:
    """슬라이드 title placeholder 우선, 없으면 첫 번째 텍스트 shape 의 첫 줄."""
    try:
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame:
            text = (title_shape.text_frame.text or "").strip()
            if text:
                return text.splitlines()[0].strip()
    except Exception:  # noqa: BLE001 — 일부 레이아웃은 .title 미존재
        pass

    return _first_text_in_shapes(slide.shapes)


def _first_text_in_shapes(shapes_iter) -> str | None:
    """shape 컬렉션에서 첫 텍스트의 첫 줄 (GroupShape 재귀)."""
    for shape in shapes_iter:
        if hasattr(shape, "shapes"):
            nested = _first_text_in_shapes(shape.shapes)
            if nested:
                return nested
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = (shape.text_frame.text or "").strip()
        except Exception:  # noqa: BLE001
            continue
        if text:
            return text.splitlines()[0].strip()
    return None


def _extract_slide_text(slide, *, warnings: list[str]) -> list[str]:
    """슬라이드 내부 모든 텍스트(텍스트 박스 + 표 + GroupShape 재귀) 추출.

    - text_frame: 단락 단위 join (paragraph 사이 줄바꿈 보존)
    - table: ` | ` separator + 행 단위 줄바꿈 (DocxParser 패턴 재사용)
    - GroupShape: `.shapes` 재귀 — 디자인 PPT 의 흔한 그룹 구조
    - picture / chart 의 캡션: 무시 (Vision 어댑터 스코프, 후속)
    """
    parts: list[str] = []
    _walk_shapes(slide.shapes, parts=parts, warnings=warnings)
    return parts


def _walk_shapes(shapes_iter, *, parts: list[str], warnings: list[str]) -> None:
    """shape 컬렉션을 재귀 순회 — GroupShape 내부 텍스트 박스도 회수."""
    for shape in shapes_iter:
        try:
            # GroupShape — 자식 shape 재귀. msopptx 의 MSO_SHAPE_TYPE.GROUP = 6 비교는
            # 라이브러리 import 회피 위해 .shapes 속성 존재 여부로 duck-typing.
            if hasattr(shape, "shapes"):
                _walk_shapes(shape.shapes, parts=parts, warnings=warnings)
                continue
            if getattr(shape, "has_text_frame", False):
                tf_text = (shape.text_frame.text or "").strip()
                if tf_text:
                    parts.append(tf_text)
            elif getattr(shape, "has_table", False):
                table_text = _table_to_text(shape.table)
                if table_text:
                    parts.append(table_text)
        except Exception as exc:  # noqa: BLE001 — shape 단위 실패 허용
            warnings.append(f"PPTX shape 추출 실패: {exc}")
            continue


def _table_to_text(table) -> str:
    """표를 chunk_filter table_noise 룰이 마킹 가능한 형태로 텍스트화.

    DocxParser `_table_to_text` 동일 패턴 — 각 행 `cell1 | cell2 | ...`,
    행 사이 `\\n`. 빈 셀은 공란.
    """
    rows_text: list[str] = []
    for row in table.rows:
        try:
            cells = [(c.text or "").strip() for c in row.cells]
            if any(cells):
                rows_text.append(" | ".join(cells))
        except Exception:  # noqa: BLE001 — 행 단위 실패 허용
            continue
    return "\n".join(rows_text)


def _vision_ocr_largest_picture(
    slide,
    *,
    slide_idx: int,
    file_name: str,
    image_parser,
    warnings: list[str],
) -> str | None:
    """슬라이드 내 가장 큰 Picture 의 image_bytes → ImageParser.parse() 결과 OCR 텍스트.

    "가장 큰" 기준 — `width * height` (pptx EMU 단위, 비교용으로 충분).
    GroupShape 내부 Picture 까지 재귀 수집.

    반환: OCR 텍스트 (caption section + ocr section 의 raw_text join). Picture 없거나
    Vision 호출 실패 시 None — caller 가 빈 슬라이드로 처리.
    """
    pictures = _collect_pictures(slide.shapes)
    if not pictures:
        return None

    # 가장 큰 picture 1장만 — Vision 비용·시간 cap
    largest = max(pictures, key=lambda p: _picture_area(p))
    try:
        image = largest.image  # python-pptx Image namedtuple-like
        blob = image.blob
        ext = (image.ext or "png").lower()
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"PPTX slide {slide_idx + 1} picture blob 추출 실패: {exc}")
        return None

    pseudo_name = f"{file_name}#slide{slide_idx + 1}.{ext}"
    try:
        result = image_parser.parse(blob, file_name=pseudo_name)
    except Exception as exc:  # noqa: BLE001 — Vision API 실패 graceful
        warnings.append(
            f"PPTX slide {slide_idx + 1} Vision OCR 실패 (graceful): {exc}"
        )
        logger.warning(
            "PPTX Vision OCR 실패 (file=%s slide=%d): %s",
            file_name, slide_idx + 1, exc,
        )
        return None

    text = (result.raw_text or "").strip()
    return text or None


def _collect_pictures(shapes_iter) -> list:
    """shape 컬렉션에서 Picture shape 만 재귀 수집 (GroupShape 내부 포함)."""
    out: list = []
    for shape in shapes_iter:
        if hasattr(shape, "shapes"):
            out.extend(_collect_pictures(shape.shapes))
            continue
        # MSO_SHAPE_TYPE.PICTURE = 13 — 라이브러리 import 회피 위해 .image 속성 duck-typing
        if hasattr(shape, "image"):
            out.append(shape)
    return out


def _picture_area(shape) -> int:
    """`width * height` (EMU 단위). 일부 shape은 None 가능 → 0 fallback."""
    w = getattr(shape, "width", None) or 0
    h = getattr(shape, "height", None) or 0
    try:
        return int(w) * int(h)
    except Exception:  # noqa: BLE001
        return 0
