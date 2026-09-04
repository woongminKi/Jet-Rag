"""S4 DOCX/PPTX 텍스트 추출 검증용 fixture 생성기.

저장소의 유일한 실제 PPTX(`브랜딩_스튜디오앤드오어.pptx`)는 **텍스트가 0자**다 —
11슬라이드 전부 이미지고 `<a:t>` 가 하나도 없다(media 106개). 운영에서는 Vision OCR 로
처리되므로 이 자산으로는 "텍스트 추출이 되는가"를 판정할 수 없다.

그래서 python-pptx 로 텍스트가 있는 fixture 를 만든다. 현행 `PptxParser` 가 실제로 다루는
구조를 전부 담는다:
  - title placeholder (section_title 의 1순위 소스)
  - 본문 텍스트 프레임 (여러 단락)
  - **group shape** (`.shapes` 재귀 — 디자인 PPT 에 흔하다)
  - table (cell 텍스트)
  - 한국어 + 영문 혼용

DOCX 쪽도 같은 이유로 합성한다. 실자산 `승인글 템플릿*.docx` 는 `.gitignore` 의 `/*.docx`
대상이라 저장소에 없다 — 그 본문을 기준선 JSON 으로 커밋하면 자산이 새는 것과 같다.
fixture 는 현행 `DocxParser` 가 실제로 쓰는 요소를 담는다:
  - `Title` / `Heading 1` / `Heading 2` 스타일 (style 이름 정규식 → section_title sticky)
  - 본문 단락, 표(`cell1 | cell2` 규칙)
  - 인라인 패턴 heading (`제 3 조` — 스타일 없이 텍스트만으로 판정되는 경로)
  - **비BMP 문자(📌)** — Python `len()`(코드포인트)과 JS `.length`(UTF-16) 차이를 회귀로 고정

사용:
    api/.venv/bin/python api/scripts/spike_ooxml_fixture.py
산출물: api/scripts/fixtures/spike_sample.pptx · spike_sample.docx
"""

from __future__ import annotations

import os

import docx as python_docx
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, "fixtures")
OUT_PPTX = os.path.join(OUT_DIR, "spike_sample.pptx")
OUT_DOCX = os.path.join(OUT_DIR, "spike_sample.docx")


def build_pptx() -> None:
    prs = Presentation()

    # --- slide 1: title + subtitle ---
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "제품 요구사항 정의서"
    s1.placeholders[1].text = "Jet-Rag 이관 스파이크 · 2026년 9월"

    # --- slide 2: title + 본문 여러 단락 ---
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "배경과 문제 정의"
    body = s2.placeholders[1].text_frame
    body.text = "현행 백엔드는 Railway 에서 구동된다."
    for line in [
        "Supabase Edge Functions 는 요청당 CPU 2초 제한이 있다.",
        "인제스트는 페이지 단위로 팬아웃해야 한다.",
        "Parser layer must run as WebAssembly.",
    ]:
        body.add_paragraph().text = line

    # --- slide 3: group shape 안의 텍스트 박스 2개 ---
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "그룹 도형 처리"
    box1 = s3.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    box1.text_frame.text = "그룹 안 첫 번째 상자"
    box2 = s3.shapes.add_textbox(Inches(1), Inches(3.2), Inches(3), Inches(1))
    box2.text_frame.text = "그룹 안 두 번째 상자 — 재귀 순회 확인용"
    # python-pptx 는 group_shapes API 가 제한적이라 XML 로 직접 묶는다.
    group = s3.shapes.add_group_shape()
    group._element.append(box1._element)
    group._element.append(box2._element)
    # group 자체의 좌표(빈 그룹으로 생성되므로 명시)
    group.left, group.top, group.width, group.height = Inches(1), Inches(2), Inches(3), Inches(2.2)

    # --- slide 4: table ---
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    s4.shapes.title.text = "판정 기준표"
    rows, cols = 3, 3
    table = s4.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(2)).table
    data = [
        ["항목", "기준", "결과"],
        ["CPU", "2초 미만", "통과"],
        ["텍스트 유사도", "0.95 이상", "통과"],
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

    prs.save(OUT_PPTX)
    print(f"→ {OUT_PPTX} ({os.path.getsize(OUT_PPTX) / 1024:.0f}KB)")


def build_docx() -> None:
    doc = python_docx.Document()

    doc.add_paragraph("이관 스파이크 문서", style="Title")
    doc.add_paragraph("표지 아래 본문 단락이다. 📌 비BMP 문자를 일부러 넣는다.")

    doc.add_paragraph("제1장 총칙", style="Heading 1")
    doc.add_paragraph("이 장은 스타일 이름으로 heading 이 판정되는 경로를 덮는다.")
    doc.add_paragraph("두 번째 본문 단락. 앞 heading 이 sticky 로 붙어야 한다.")

    doc.add_paragraph("제1절 목적", style="Heading 2")
    doc.add_paragraph("Heading 2 는 styles.xml 에 `heading 2` 로 적힌다 — 정규식이 대소문자를 무시해야 잡힌다.")

    # 스타일 없이 텍스트 패턴만으로 heading 이 되는 경로 (Normal 스타일)
    doc.add_paragraph("제 3 조 (적용범위)")
    doc.add_paragraph("인라인 패턴 heading 뒤의 본문.")

    table = doc.add_table(rows=3, cols=3)
    data = [
        ["항목", "기준", "결과"],
        ["CPU", "2초 미만", "통과"],
        ["유사도", "0.95 이상", "통과 📌"],
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    doc.add_paragraph("부칙")
    doc.add_paragraph("표 뒤 단락 — 표와 단락의 XML 순서가 보존되는지 확인용.")

    doc.save(OUT_DOCX)
    print(f"→ {OUT_DOCX} ({os.path.getsize(OUT_DOCX) / 1024:.0f}KB)")


def main() -> None:
    os.makedirs(OUT_DIR, exist_ok=True)
    build_pptx()
    build_docx()


if __name__ == "__main__":
    main()
